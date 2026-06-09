#!/usr/bin/env python3
# Generate a PowerPoint describing the synthesizable instruction tracer
# (instr_tracer_synth) architecture and its CVA6 co-simulation verification flow.
#
#   pip3 install --user python-pptx
#   python3 scripts/gen_tracer_arch_ppt.py
#
# Output: docs/instr_tracer_synth_arch.pptx

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------- palette
NAVY    = RGBColor(0x1F, 0x3A, 0x5F)
BLUE    = RGBColor(0x2E, 0x6D, 0xB4)
SYNTH   = RGBColor(0xDC, 0xEA, 0xFB)   # synthesizable RTL
SYNTH_B = RGBColor(0x2E, 0x6D, 0xB4)
SIM     = RGBColor(0xFC, 0xE7, 0xCE)   # simulation-only
SIM_B   = RGBColor(0xE0, 0x8A, 0x1E)
GOLD    = RGBColor(0xD9, 0xF0, 0xD3)   # golden reference
GOLD_B  = RGBColor(0x4C, 0x9A, 0x3F)
FILEC   = RGBColor(0xEA, 0xEA, 0xEA)   # output file
FILE_B  = RGBColor(0x77, 0x77, 0x77)
REDF    = RGBColor(0xFD, 0xDD, 0xD6)   # warning / overflow
REDB    = RGBColor(0xC8, 0x4A, 0x33)
GREYTX  = RGBColor(0x40, 0x40, 0x40)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MONO    = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def title_bar(s, text, sub=None):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.92))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background(); bar.shadow.inherit = False
    tf = bar.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.35)
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(25); p.font.bold = True; p.font.color.rgb = WHITE
    if sub:
        p2 = tf.add_paragraph(); p2.text = sub
        p2.font.size = Pt(12.5); p2.font.color.rgb = RGBColor(0xC4, 0xD4, 0xE6)
    acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.92), prs.slide_width, Inches(0.06))
    acc.fill.solid(); acc.fill.fore_color.rgb = BLUE
    acc.line.fill.background(); acc.shadow.inherit = False


def box(s, x, y, w, h, title, body=None, fill=SYNTH, border=SYNTH_B,
        tcolor=NAVY, tsize=13, bsize=9.5, bold=True,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, align=PP_ALIGN.CENTER, mono_body=False):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = border; sp.line.width = Pt(1.5)
    sp.shadow.inherit = False
    tf = sp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.07); tf.margin_right = Inches(0.07)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = title
    r.font.size = Pt(tsize); r.font.bold = bold; r.font.color.rgb = tcolor
    if body:
        for line in body:
            pp = tf.add_paragraph(); pp.alignment = align; pp.space_before = Pt(1)
            rr = pp.add_run(); rr.text = line
            rr.font.size = Pt(bsize); rr.font.color.rgb = GREYTX
            if mono_body:
                rr.font.name = MONO
    return sp


def arrow(s, x1, y1, x2, y2, color=GREYTX, width=2.0):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(width)
    cn.shadow.inherit = False
    ln = cn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return cn


def label(s, x, y, w, text, size=10, color=GREYTX, bold=False,
          align=PP_ALIGN.CENTER, h=0.3, italic=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    r.font.italic = italic
    return tb


def caption(s, x, y, w, text, color=BLUE):
    return label(s, x, y, w, text, size=9.5, color=color, align=PP_ALIGN.CENTER,
                 h=0.45, italic=True)


def legend(s, x, y, items):
    cx = x
    for col, brd, txt in items:
        sw = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(y),
                                Inches(0.26), Inches(0.26))
        sw.fill.solid(); sw.fill.fore_color.rgb = col
        sw.line.color.rgb = brd; sw.line.width = Pt(1.25); sw.shadow.inherit = False
        label(s, cx + 0.32, y - 0.02, 2.6, txt, size=10, align=PP_ALIGN.LEFT)
        cx += 0.34 + 0.105 * len(txt) + 0.25


def badge(s, x, y, n, color=BLUE, d=0.42):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = color
    c.line.color.rgb = WHITE; c.line.width = Pt(1.0); c.shadow.inherit = False
    tf = c.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE
    return c


def step_head(s, x, y, n, text, w=3.0):
    badge(s, x, y, n)
    label(s, x + 0.52, y + 0.04, w - 0.52, text, size=12, color=NAVY,
          bold=True, align=PP_ALIGN.LEFT, h=0.4)


# =========================================================================== SLIDE 1
s = slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background(); bg.shadow.inherit = False
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), prs.slide_width, Inches(0.06))
band.fill.solid(); band.fill.fore_color.rgb = BLUE; band.line.fill.background(); band.shadow.inherit = False
label(s, 0.8, 1.35, 11.7, "Synthesizable Instruction Tracer for CVA6",
      size=40, color=WHITE, bold=True, align=PP_ALIGN.LEFT, h=1.1)
label(s, 0.82, 2.75, 11.7, "Architecture  &  Spike Co-Simulation Verification Flow",
      size=20, color=RGBColor(0xBF, 0xD2, 0xE8), align=PP_ALIGN.LEFT, h=0.6)
label(s, 0.82, 3.7, 11.7,
      "instr_tracer_synth  —  a hardware-implementable re-creation of CVA6's "
      "simulation-only instr_tracer\nthat reproduces the exact riscv::spikeCommitLog() "
      "commit log using only synthesizable RTL.",
      size=14.5, color=RGBColor(0xD8, 0xE2, 0xEF), align=PP_ALIGN.LEFT, h=1.0)
label(s, 0.82, 6.55, 11.7, "hardware/src/trace/  ·  ara  ·  CVA6 + Ara vector unit",
      size=12, color=RGBColor(0x90, 0xA6, 0xC0), align=PP_ALIGN.LEFT)

# =========================================================================== SLIDE 2
s = slide()
title_bar(s, "Concept: split the problem on-chip / off-chip",
          "Hardware cannot synthesize ASCII strings or file I/O — so the chip captures + packs; the host decodes.")
box(s, 0.7, 1.6, 4.6, 1.7, "ON-CHIP  (synthesizable)",
    ["Capture retiring instructions at the", "commit stage → pack into a fixed-",
     "layout binary record + stream out", "(instr_tracer_synth.sv)"],
    fill=SYNTH, border=SYNTH_B, tsize=15)
box(s, 8.0, 1.6, 4.6, 1.7, "OFF-CHIP  (host)",
    ["Decode the binary stream into", "Spike-format ASCII commit log",
     "(scripts/spike_trace_decode.py,", "or the SystemVerilog sim sink)"],
    fill=SIM, border=SIM_B, tsize=15)
arrow(s, 5.3, 2.45, 8.0, 2.45, color=BLUE, width=2.75)
label(s, 5.2, 2.05, 2.9, "binary records", size=11, color=BLUE, bold=True)

label(s, 0.7, 3.75, 12.0, "Non-synthesizable construct  →  synthesizable equivalent",
      size=15, color=NAVY, bold=True, align=PP_ALIGN.LEFT)
rows = [
    ("SystemVerilog class / new",          "packed struct  commit_log_pkt_t"),
    ("dynamic queue  logic[..][$]",        "fixed-depth HW FIFO (fifo_v3 / addr_fifo)"),
    ("string + $sformatf",                 "fixed-layout binary record (300 bit @XLEN=64)"),
    ("$fopen / $fwrite / $fclose",         "ready / valid streaming port"),
    ("interface + clocking block",         "plain module ports (packed struct)"),
    ("initial / forever / task",           "always_ff / always_comb"),
    ("gp_reg_file / fp_reg_file shadow",   "synthesizable gp/fp_reg_file_q[32]"),
]
tbl = s.shapes.add_table(len(rows) + 1, 2, Inches(0.7), Inches(4.25),
                         Inches(11.9), Inches(2.7)).table
tbl.columns[0].width = Inches(5.3); tbl.columns[1].width = Inches(6.6)
for c, txt in enumerate(["simulation tracer (NOT synthesizable)", "instr_tracer_synth (synthesizable)"]):
    cell = tbl.cell(0, c); cell.text = txt
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    pr = cell.text_frame.paragraphs[0]; pr.font.size = Pt(12); pr.font.bold = True; pr.font.color.rgb = WHITE
for ri, (a, b) in enumerate(rows, start=1):
    for ci, txt in enumerate((a, b)):
        cell = tbl.cell(ri, ci); cell.text = txt
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if ri % 2 else RGBColor(0xEE, 0xF3, 0xF9)
        pr = cell.text_frame.paragraphs[0]
        pr.font.size = Pt(11); pr.font.name = MONO
        pr.font.color.rgb = BLUE if ci == 1 else GREYTX

# =========================================================================== SLIDE: PINOUT
s = slide()
title_bar(s, "Module pinout  —  every input / output pin",
          "4 input groups (left) feed the tracer; the ready/valid trace-out port (right) streams packed beats.")
_png = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                    "hardware", "src", "trace", "instr_tracer_synth_pinout.png")
_ph = 5.95
_pw = _ph * 1820.0 / 1836.0
s.shapes.add_picture(_png, Inches(0.30), Inches(1.20), height=Inches(_ph))
_rx, _rw = 6.75, 6.25
box(s, _rx, 1.25, _rw, 0.92, "CLOCK / RESET / CONTROL    (4 in)",
    ["clk_i · rst_ni · flush_i · testmode_i"],
    fill=RGBColor(0xEC, 0xEC, 0xEC), border=RGBColor(0x6A, 0x6A, 0x6A),
    tsize=12, bsize=10, mono_body=True, align=PP_ALIGN.LEFT)
box(s, _rx, 2.30, _rw, 1.28, "COMMIT-STAGE OBSERVATION    (10 in)",
    ["commit_instr_i · commit_ack_i · instr_word_i",
     "waddr_i · wdata_i · we_gpr_i · we_fpr_i",
     "priv_lvl_i · debug_mode_i · exception_i"],
    fill=SYNTH, border=SYNTH_B, tsize=12, bsize=10, mono_body=True, align=PP_ALIGN.LEFT)
box(s, _rx, 3.71, _rw, 1.00, "CSR WRITE OBSERVATION    (5 in)",
    ["csr_commit_i · csr_op_i · csr_waddr_i", "csr_operand_i · csr_old_i"],
    fill=RGBColor(0xEE, 0xE6, 0xF7), border=RGBColor(0x78, 0x46, 0xA0),
    tsize=12, bsize=10, mono_body=True, align=PP_ALIGN.LEFT)
box(s, _rx, 4.84, _rw, 1.00, "LSU MEMORY-ACCESS OBSERVATION    (9 in)",
    ["st_valid/paddr/data/size_i", "ld_valid/kill/paddr/size_i · flush_addr_i"],
    fill=RGBColor(0xFC, 0xEE, 0xDC), border=RGBColor(0xC4, 0x76, 0x1C),
    tsize=12, bsize=10, mono_body=True, align=PP_ALIGN.LEFT)
box(s, _rx, 5.97, _rw, 1.00, "STREAMING TRACE-OUT    (3 out + 1 in)",
    ["trace_valid_o · trace_beat_o · overflow_o    (out)",
     "trace_ready_i    (in, back-pressure)"],
    fill=GOLD, border=GOLD_B, tsize=12, bsize=10, mono_body=True, align=PP_ALIGN.LEFT)

# =========================================================================== SLIDE 3A
s = slide()
title_bar(s, "Internal microarchitecture  —  the beat builder & its two memories",
          "One combinational beat builder reads the live commit signals + two on-chip memories; a FIFO then streams the result out.")
_SIGF  = RGBColor(0xE9, 0xEC, 0xF1)
_MEMF  = RGBColor(0xFC, 0xEF, 0xD9)
_MEMB  = RGBColor(0xC4, 0x76, 0x1C)
_BEATF = RGBColor(0xC9, 0xDE, 0xF7)
legend(s, 0.70, 6.45, [(_SIGF, NAVY, "input"), (_MEMF, _MEMB, "memory"), (_BEATF, BLUE, "beat builder")])

# --- top : the live inputs, each dropping straight into its block ----------
box(s, 0.70, 1.18, 2.70, 0.95, "Write-back bus",
    ["waddr · wdata", "we_gpr / we_fpr"], fill=_SIGF, border=NAVY, tsize=12.5, bsize=9.5, mono_body=True)
box(s, 4.85, 1.18, 3.63, 0.95, "commit-stage signals",
    ["commit_instr · instr_word · commit_ack", "priv · debug · exception_i · csr_*"],
    fill=_SIGF, border=NAVY, tsize=12.5, bsize=9.5, mono_body=True)
box(s, 9.95, 1.18, 2.70, 0.95, "LSU addresses",
    ["st_valid/paddr/data/size", "ld_valid/kill/paddr/size"], fill=_SIGF, border=NAVY, tsize=12.5, bsize=9.5, mono_body=True)

# --- middle : two memories flanking the builder ----------------------------
box(s, 0.70, 2.55, 2.70, 1.95, "Shadow regfile",
    ["gp / fp_reg_file_q[32]", "always_ff", "",
     "recovers rd's value when", "it isn't on the write-back", "bus this cycle"],
    fill=_MEMF, border=_MEMB, tsize=14, bsize=9.5)
box(s, 5.00, 2.45, 3.33, 2.25, "BEAT  BUILDER",
    ["always_comb", "",
     "per-port packet:", "priv·pc·instr·rd·we·wdata", "compressed · retired · mem",
     "", "port-0 overrides:  exception · CSR", "",
     "→  commit_log_beat_t"], fill=_BEATF, border=BLUE, tsize=15, bsize=9.5)
box(s, 9.95, 2.55, 2.70, 1.95, "Address FIFOs",
    ["store + load · depth 16", "+ byte_ror64", "",
     "realign LSU addresses", "to commit order"], fill=_MEMF, border=_MEMB, tsize=14, bsize=9.5)

# inputs drop straight down into their block
arrow(s, 2.05, 2.13, 2.05, 2.55)
arrow(s, 6.665, 2.13, 6.665, 2.45, color=BLUE, width=2.4)
arrow(s, 11.30, 2.13, 11.30, 2.55)

# the builder reads the two memories; pops the address FIFO at commit
arrow(s, 3.40, 3.35, 5.00, 3.35, color=_MEMB, width=2.2)        # shadow -> beat
label(s, 3.40, 3.02, 1.60, "reg[rd]", size=10, color=_MEMB, bold=True)
arrow(s, 9.95, 3.35, 8.33, 3.35, color=_MEMB, width=2.2)        # addr FIFO -> beat (mem)
label(s, 8.33, 3.02, 1.62, "mem addr / data / size", size=9, color=_MEMB, bold=True)
arrow(s, 8.33, 4.05, 9.95, 4.05, color=GREYTX, width=1.75)      # beat -> addr FIFO (pop)
label(s, 8.33, 4.08, 1.62, "pop @ commit", size=9, color=GREYTX)

# --- bottom : straight backbone out ----------------------------------------
box(s, 5.00, 5.00, 3.33, 0.85, "Trace FIFO",
    ["fifo_v3 · depth 32   ·   push = |valid & ~full"], fill=SYNTH, border=SYNTH_B, tsize=14, bsize=9.5)
box(s, 5.00, 6.10, 3.33, 0.85, "Trace-out port",
    ["ready / valid   ·   trace_valid_o · trace_beat_o · trace_ready_i"], fill=SYNTH, border=SYNTH_B, tsize=14, bsize=8.5)
box(s, 9.20, 5.00, 2.80, 0.85, "overflow_o",
    ["FIFO full → trace gap"], fill=REDF, border=REDB, tsize=13, bsize=9.5)

arrow(s, 6.665, 4.70, 6.665, 5.00, color=BLUE, width=2.6)       # beat -> trace FIFO
arrow(s, 6.665, 5.85, 6.665, 6.10, color=BLUE, width=2.6)       # trace FIFO -> out
arrow(s, 8.33, 5.42, 9.20, 5.42, color=REDB, width=1.75)        # trace FIFO -> overflow
label(s, 8.33, 5.10, 0.87, "full", size=9, color=REDB, bold=True)

# =========================================================================== SLIDE: DATAPATH DETAIL
s = slide()
title_bar(s, "Detailed datapath  —  sub-blocks & signals",
          "Same flanked-builder view, zoomed in: value recovery on the left, memory re-alignment on the right, block-by-block.")
_SRC   = RGBColor(0xEC, 0xEE, 0xF2)
_ORN   = RGBColor(0xC4, 0x76, 0x1C)
_MEMF  = RGBColor(0xFC, 0xEF, 0xD9)
_BEATF = RGBColor(0xC9, 0xDE, 0xF7)
legend(s, 0.30, 1.28, [(_SRC, NAVY, "input"), (_MEMF, _ORN, "memory"),
                       (SYNTH, SYNTH_B, "logic"), (_BEATF, BLUE, "beat builder")])

# --- centre : the beat builder ---------------------------------------------
bbx, bby, bbw, bbh = 5.55, 2.10, 2.25, 2.80
box(s, bbx, bby, bbw, bbh, "BEAT  BUILDER",
    ["always_comb", "", "builds beat.pkt[p]", "for both ports", "+ port-0",
     "exception / CSR", "", "drives", "st_pop / ld_pop", "", "beat.valid =", "  commit_ack"],
    fill=_BEATF, border=BLUE, bsize=8.5, tsize=13)

# --- top : commit-stage signals drop straight in ---------------------------
box(s, 5.20, 1.10, 2.95, 0.82, "commit-stage signals",
    ["commit_instr · instr_word · priv · debug", "exception_i · csr_*   (→ port 0)"],
    fill=_SRC, border=NAVY, tsize=10.5, bsize=7.8, mono_body=True)
arrow(s, 6.675, 1.92, 6.675, 2.10, color=BLUE, width=2.2)

# --- left flank : value recovery (flows → into the builder) ----------------
box(s, 0.30, 2.55, 1.55, 0.95, "Write-back bus",
    ["we_gpr_i / we_fpr_i", "waddr_i · wdata_i"], fill=_SRC, border=NAVY, tsize=9.5, bsize=7.6, mono_body=True)
box(s, 1.98, 2.55, 1.65, 0.95, "Shadow GP/FP regfile",
    ["gp/fp_reg_file_q[32]", "always_ff"], fill=_MEMF, border=_ORN, tsize=9.5, bsize=7.8)
box(s, 3.78, 2.50, 1.42, 1.05, "Result MUX",
    ["we ? wdata_i", " : rd_fpr ? fp[rd]", " :          gp[rd]"], tsize=9.5, bsize=7.6, mono_body=True)
arrow(s, 1.85, 3.02, 1.98, 3.02)
arrow(s, 3.63, 3.02, 3.78, 3.02)
label(s, 3.13, 2.66, 1.0, "reg[rd]", size=7.5, color=_ORN)
arrow(s, 5.20, 3.02, bbx, 3.05, color=BLUE, width=2.0)
label(s, 4.85, 2.64, 1.0, "result", size=8, color=BLUE)
caption(s, 0.30, 3.62, 5.0, "→ recovers rd's committed value when it is not on the write-back bus this cycle")

# --- right flank : memory re-alignment (flows ← into the builder) ----------
box(s, 8.15, 2.55, 1.35, 0.95, "byte_ror64",
    ["undo data_align", "(rotate right)"], tsize=9.5, bsize=7.6)
box(s, 9.65, 2.55, 1.70, 0.95, "Store addr FIFO",
    ["addr_fifo · depth 16", "pop 0 / 1"], fill=_MEMF, border=_ORN, tsize=9, bsize=7.8)
box(s, 11.50, 2.55, 1.50, 0.95, "Store gen (LSU)",
    ["st_valid_i", "st_paddr/data/size_i"], fill=_SRC, border=NAVY, tsize=9, bsize=7.4, mono_body=True)
arrow(s, 11.50, 3.02, 11.35, 3.02)
arrow(s, 9.65, 3.02, 9.50, 3.02)
label(s, 9.10, 2.66, 1.0, "dout0.data", size=7, color=_ORN)
arrow(s, 8.15, 3.05, bbx + bbw, 3.05, color=_ORN, width=2.0)
label(s, 7.55, 2.64, 1.0, "mem", size=8, color=_ORN)

box(s, 9.65, 3.75, 1.70, 0.95, "Load addr FIFO",
    ["addr_fifo · depth 16", "pop 0..2 · data0/1"], fill=_MEMF, border=_ORN, tsize=9, bsize=7.6)
box(s, 11.50, 3.75, 1.50, 0.95, "Load gen (LSU)",
    ["ld_valid_i & ~ld_kill_i", "ld_paddr/size_i"], fill=_SRC, border=NAVY, tsize=9, bsize=7.2, mono_body=True)
arrow(s, 11.50, 4.22, 11.35, 4.22)
arrow(s, 9.65, 4.22, bbx + bbw, 4.20, color=_ORN, width=2.0)
label(s, 7.95, 3.82, 1.7, "ld_dout0/1 → mem_addr", size=7.5, color=_ORN)
caption(s, 7.90, 4.80, 5.2, "→ FIFOs realign out-of-order LSU addresses; popped at commit by st_pop / ld_pop (0..2)")

# --- bottom : straight backbone out ----------------------------------------
box(s, 5.30, 5.25, 2.75, 0.85, "Trace FIFO",
    ["fifo_v3 · depth 32", "push=|valid&~full · pop=valid&ready"], fill=SYNTH, border=SYNTH_B, tsize=12, bsize=8)
box(s, 5.30, 6.35, 2.75, 0.85, "Trace-out port",
    ["trace_valid_o=~empty (out) · trace_beat_o (out)", "trace_ready_i (in)"],
    fill=SYNTH, border=SYNTH_B, tsize=12, bsize=8)
box(s, 8.60, 5.25, 2.70, 0.85, "overflow_o   (out)",
    ["(any_valid & full)", "| st/ld of | uf"], fill=REDF, border=REDB, tsize=11.5, bsize=8.5)
arrow(s, 6.675, 4.90, 6.675, 5.25, color=BLUE, width=2.6)
arrow(s, 6.675, 6.10, 6.675, 6.35, color=BLUE, width=2.6)
arrow(s, 8.05, 5.67, 8.60, 5.67, color=REDB, width=1.75)
label(s, 8.05, 5.36, 0.55, "full", size=8.5, color=REDB, bold=True)

# =========================================================================== SLIDE 3B
s = slide()
title_bar(s, "What goes into one trace record  &  the 3 packing rules",
          "Each retired instruction becomes a packed commit_log_pkt_t; two of them form a beat (program order kept).")

# left: record format
label(s, 0.55, 1.20, 5.9, "One beat  =  commit_log_beat_t", size=15, color=NAVY,
      bold=True, align=PP_ALIGN.LEFT)
box(s, 0.55, 1.65, 5.75, 0.78, "commit_log_beat_t",
    ["valid[NR_COMMIT_PORTS]   +   2 × commit_log_pkt_t   (port 0 = older)"],
    fill=RGBColor(0xE9, 0xEC, 0xF1), border=NAVY, tsize=12.5, bsize=9.5, mono_body=True)
arrow(s, 3.4, 2.43, 3.4, 2.72, color=NAVY)

fields = [
    ("Identity", "priv · debug · retired · compressed · rd (+rd_fpr) · pc · instr"),
    ("Result",   "we · wdata"),
    ("Exception","ex_valid · cause · tval"),
    ("Memory",   "mem_op · mem_addr · mem_data · mem_size"),
    ("CSR write","csr_we · csr_addr · csr_wdata"),
]
ft = s.shapes.add_table(len(fields) + 1, 2, Inches(0.55), Inches(2.78),
                        Inches(5.75), Inches(3.0)).table
ft.columns[0].width = Inches(1.5); ft.columns[1].width = Inches(4.25)
for c, txt in enumerate(["group", "fields of commit_log_pkt_t"]):
    cell = ft.cell(0, c); cell.text = txt
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    pr = cell.text_frame.paragraphs[0]; pr.font.size = Pt(11.5); pr.font.bold = True; pr.font.color.rgb = WHITE
for ri, (g, f) in enumerate(fields, start=1):
    cg = ft.cell(ri, 0); cg.text = g
    cg.fill.solid(); cg.fill.fore_color.rgb = SYNTH
    pg = cg.text_frame.paragraphs[0]; pg.font.size = Pt(11); pg.font.bold = True; pg.font.color.rgb = NAVY
    cf = ft.cell(ri, 1); cf.text = f
    cf.fill.solid(); cf.fill.fore_color.rgb = WHITE if ri % 2 else RGBColor(0xEE, 0xF3, 0xF9)
    pf = cf.text_frame.paragraphs[0]; pf.font.size = Pt(10.5); pf.font.name = MONO; pf.font.color.rgb = GREYTX
label(s, 0.55, 5.95, 5.75, "300 bits @ XLEN=64  ·  packed MSB-first in declaration order",
      size=10.5, color=GREYTX, align=PP_ALIGN.LEFT, italic=True)

# right: 3 rules
label(s, 6.7, 1.20, 6.1, "How each field is filled  —  3 rules", size=15, color=NAVY,
      bold=True, align=PP_ALIGN.LEFT)
box(s, 6.7, 1.68, 6.1, 1.30, "Rule 1   Result value  (a 3-way mux)",
    ["we_gpr | we_fpr  →  wdata_i   (written this cycle)",
     "else rd is FPR   →  fp_reg_file_q[rd]",
     "else             →  gp_reg_file_q[rd]"],
    fill=SYNTH, border=SYNTH_B, tsize=12.5, bsize=9.7, mono_body=True, align=PP_ALIGN.LEFT)
box(s, 6.7, 3.12, 6.1, 1.55, "Rule 2   Memory token  ( “mem 0x..” )",
    ["commit & STORE (non-AMO) → MEM_STORE, pop store FIFO,",
     "                          data via byte_ror64",
     "commit & LOAD            → MEM_LOAD (up to 2/cycle),",
     "                          addr from load FIFO",
     "otherwise                → MEM_NONE"],
    fill=SYNTH, border=SYNTH_B, tsize=12.5, bsize=9.5, mono_body=True, align=PP_ALIGN.LEFT)
box(s, 6.7, 4.82, 6.1, 1.55, "Rule 3   Port-0-only overrides",
    ["exception.valid & !(debug & BREAKPOINT)",
     "        → ex_valid, cause, tval",
     "csr_commit & op ∈ {WRITE, SET, CLEAR}",
     "        → csr_we, reconstruct csr_wdata",
     "(both always reported on the OLDER port 0)"],
    fill=SYNTH, border=SYNTH_B, tsize=12.5, bsize=9.5, mono_body=True, align=PP_ALIGN.LEFT)

# =========================================================================== SLIDE: BIT LAYOUT
s = slide()
title_bar(s, "commit_log_pkt_t  —  packed bit layout (MSB-first)",
          "One retired-instruction record. A valid[2] mask + two records form a commit_log_beat_t.")

def _fieldtbl(x, fields):
    t = s.shapes.add_table(len(fields) + 1, 3, Inches(x), Inches(1.55),
                           Inches(5.55), Inches(0.34 * (len(fields) + 1))).table
    t.columns[0].width = Inches(1.55); t.columns[1].width = Inches(0.85); t.columns[2].width = Inches(3.15)
    for c, hh in enumerate(["field", "bits", "meaning"]):
        cell = t.cell(0, c); cell.text = hh
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        pr = cell.text_frame.paragraphs[0]; pr.font.size = Pt(11); pr.font.bold = True; pr.font.color.rgb = WHITE
    for ri, (fn, bits, mean, grp) in enumerate(fields, start=1):
        for ci, v in enumerate((fn, str(bits), mean)):
            cell = t.cell(ri, ci); cell.text = v
            cell.fill.solid(); cell.fill.fore_color.rgb = grp
            pr = cell.text_frame.paragraphs[0]
            pr.font.size = Pt(10.5 if ci != 2 else 9.5)
            pr.font.color.rgb = NAVY if ci == 0 else GREYTX
            if ci == 0:
                pr.font.name = MONO; pr.font.bold = True
    return t

_GI = RGBColor(0xDC, 0xEA, 0xFB); _GR = RGBColor(0xDF, 0xEF, 0xE0); _GE = RGBColor(0xFB, 0xE5, 0xE0)
_GM = RGBColor(0xFC, 0xEE, 0xDC); _GC = RGBColor(0xEE, 0xE6, 0xF7)
_fieldtbl(0.45, [
    ("priv", 2, "privilege  M / S / U", _GI), ("debug", 1, "debug-mode flag", _GI),
    ("ex_valid", 1, "exception present", _GE), ("retired", 1, "commit_ack → emit line", _GI),
    ("compressed", 1, "RVC 16-bit instr", _GI), ("we", 1, "wrote a register", _GR),
    ("rd_fpr", 1, "dest is FP register", _GI), ("rd", 5, "dest register index", _GI),
    ("pc", 64, "program counter (XLEN)", _GI), ("instr", 32, "raw instruction word", _GI)])
_fieldtbl(7.05, [
    ("wdata", 64, "committed result (XLEN)", _GR), ("cause", 64, "exception cause (XLEN)", _GE),
    ("tval", 64, "trap value (XLEN)", _GE), ("mem_op", 2, "NONE / LOAD / STORE", _GM),
    ("mem_addr", 64, "physical address", _GM), ("mem_data", 64, "value loaded / stored", _GM),
    ("mem_size", 2, "B / H / W / D", _GM), ("csr_we", 1, "a CSR was written", _GC),
    ("csr_addr", 12, "CSR address", _GC), ("csr_wdata", 64, "CSR value (pre-WARL, XLEN)", _GC)])
box(s, 0.45, 6.35, 5.55, 0.78, "record width = 510 bits  @ XLEN=64",
    ["Spike core fields priv…tval = 301  ·  memory +132  ·  CSR +77"],
    fill=SYNTH, border=SYNTH_B, tsize=13, bsize=9.5)
box(s, 7.05, 6.35, 5.55, 0.78, "commit_log_beat_t = valid[2] + 2 × pkt",
    ["= 1022 bits  ·  port 0 = older instruction (program order)"],
    fill=RGBColor(0xE9, 0xEC, 0xF1), border=NAVY, tsize=13, bsize=9.5)

# =========================================================================== SLIDE: ADDR FIFO
s = slide()
title_bar(s, "Address-FIFO re-alignment  (instr_tracer_addr_fifo)",
          "The LSU generates addresses out of program order; bounded FIFOs realign each to the retiring instruction.")
_ORN = RGBColor(0xC4, 0x76, 0x1C)
box(s, 0.5, 1.55, 2.5, 1.25, "LSU",
    ["store_unit / load_unit", "generates paddr", "OUT of program order"],
    fill=RGBColor(0xEC, 0xEE, 0xF2), border=NAVY, tsize=14, bsize=9.5)
label(s, 4.30, 1.42, 4.6, "store / load address FIFO   (depth 16)", size=11, color=NAVY, bold=True, align=PP_ALIGN.LEFT)
for _i in range(5):
    box(s, 4.30 + _i * 0.92, 1.78, 0.88, 0.85, "", None, fill=SYNTH, border=SYNTH_B)
label(s, 4.30, 2.68, 0.88, "head", size=9, color=BLUE)
label(s, 4.30 + 4 * 0.92, 2.68, 0.88, "tail", size=9, color=BLUE)
box(s, 10.20, 1.55, 2.60, 1.25, "Beat builder",
    ["pops the oldest entry", "for the retiring", "load / store"], fill=SYNTH, border=SYNTH_B, tsize=14, bsize=9.5)
arrow(s, 3.00, 2.18, 4.28, 2.18, color=_ORN, width=2.4)
label(s, 3.00, 1.80, 1.3, "push", size=9, color=_ORN, bold=True)
arrow(s, 8.95, 2.18, 10.20, 2.18, color=_ORN, width=2.4)
label(s, 8.95, 1.80, 1.3, "pop @ commit", size=9, color=_ORN, bold=True, align=PP_ALIGN.LEFT)
box(s, 0.5, 3.25, 12.3, 3.45, "How it works",
    ["• push  —  store: st_valid_i pushes {paddr, data, size};   load: ld_valid_i & ~ld_kill_i pushes {paddr, size}",
     "• pop @ commit  —  STORE (non-AMO) pops 1;   LOAD pops 0..2 (two loads may retire per cycle) → exposes data0_o & data1_o",
     "• AMOs report fu == STORE but never push the store buffer → they must NOT pop  (queue stays aligned)",
     "• byte_ror64 undoes CVA6 data_align (rotate-left by addr[2:0]) so mem_data carries the natural stored value",
     "• flush_addr_i clears the queue on a pipeline flush (branch mispredict / exception)",
     "• depth = 2^AW = 16, a power of two → natural pointer wrap-around (no occupancy compare needed)",
     "• over / underflow → overflow_o : a mem token may be misaligned, trace no longer address-accurate"],
    fill=RGBColor(0xF3, 0xF6, 0xFA), border=NAVY, tsize=15, bsize=11.5, align=PP_ALIGN.LEFT)

# =========================================================================== SLIDE: EXAMPLE
s = slide()
title_bar(s, "Worked example  —  anatomy of a Spike commit-log line",
          "How packed fields become Spike-format ASCII (printed by the sim sink or the host decoder).")
_cb = box(s, 0.5, 1.30, 12.3, 1.05, "", None, fill=RGBColor(0x1B, 0x20, 0x28),
          border=RGBColor(0x1B, 0x20, 0x28), shape=MSO_SHAPE.RECTANGLE)
_tf = _cb.text_frame; _tf.vertical_anchor = MSO_ANCHOR.MIDDLE; _tf.margin_left = Inches(0.2)
for _i, _e in enumerate([
        "0 0x0000000080002a40 (0x0007a283) x 5 0x000000000000002a mem 0x0000000080004000",
        "3 0x0000000080000100 (0x30529073) c773_mtvec 0x0000000080000200"]):
    _p = _tf.paragraphs[0] if _i == 0 else _tf.add_paragraph()
    _r = _p.add_run(); _r.text = _e; _r.font.size = Pt(12.5); _r.font.name = MONO
    _r.font.color.rgb = RGBColor(0x8C, 0xE0, 0x9A)
_toks = [
    ("0  /  3", "priv", "decimal: U=0, S=1, M=3"),
    ("0x…2a40", "pc", "0x + 16 hex digits (zero-padded)"),
    ("(0x0007a283)", "instr", "instr[1:0]=2'b11 → 8-hex; RVC → (0x____)"),
    ("x 5 0x…2a", "rd · we · wdata", "'x'/'f'+idx; printed when rd_fpr or rd≠0"),
    ("mem 0x…4000", "mem_op = LOAD", "load: addr only (value shown as the rd write)"),
    ("c773_mtvec 0x…200", "csr_we/addr/wdata", "c<dec-addr>_<name>;  name via csr_name()"),
]
_t = s.shapes.add_table(len(_toks) + 1, 3, Inches(0.5), Inches(2.65), Inches(7.6), Inches(3.0)).table
_t.columns[0].width = Inches(2.55); _t.columns[1].width = Inches(1.95); _t.columns[2].width = Inches(3.10)
for _c, _h in enumerate(["token", "← field(s)", "rule"]):
    _cell = _t.cell(0, _c); _cell.text = _h
    _cell.fill.solid(); _cell.fill.fore_color.rgb = NAVY
    _pr = _cell.text_frame.paragraphs[0]; _pr.font.size = Pt(11); _pr.font.bold = True; _pr.font.color.rgb = WHITE
for _ri, (_a, _b, _cc) in enumerate(_toks, start=1):
    for _ci, _v in enumerate((_a, _b, _cc)):
        _cell = _t.cell(_ri, _ci); _cell.text = _v
        _cell.fill.solid(); _cell.fill.fore_color.rgb = WHITE if _ri % 2 else RGBColor(0xEE, 0xF3, 0xF9)
        _pr = _cell.text_frame.paragraphs[0]; _pr.font.size = Pt(10 if _ci != 2 else 9.5)
        _pr.font.color.rgb = NAVY if _ci == 0 else GREYTX
        if _ci == 0:
            _pr.font.name = MONO
box(s, 8.4, 2.65, 4.45, 3.0, "Formatting rules",
    ["• line emitted when retired (commit_ack)",
     "   & !debug_mode — NOT gated by exceptions",
     "• priv printed in decimal (M=3, S=1, U=0)",
     "• pc / values: 0x + 16 hex digits",
     "   (spikeCommitLog hardcodes logic[63:0])",
     "• RVC if instr[1:0]≠2'b11 → (0x____)",
     "• STORE → mem 0x<addr> 0x<data>;",
     "   LOAD → mem 0x<addr>",
     "• reg spacing: Spike 'x 5' vs ours 'x5' —",
     "   compare_spike_log.py normalizes it"],
    fill=RGBColor(0xF3, 0xF6, 0xFA), border=NAVY, tsize=13, bsize=10, align=PP_ALIGN.LEFT)

# =========================================================================== SLIDE 4
s = slide()
title_bar(s, "Integration into CVA6  —  non-invasive  bind",
          "instr_tracer_synth_bind.sv attaches the tracer to every `ariane` core without editing the CVA6 submodule.")
legend(s, 0.7, 1.12, [(SYNTH, SYNTH_B, "synthesizable"), (SIM, SIM_B, "simulation-only"),
                      (FILEC, FILE_B, "output file")])
box(s, 1.6, 1.65, 10.1, 1.55, "module  ariane   (CVA6 core)",
    ["IF → ID → Issue → EX → Commit",
     "commit_instr_id_commit  commit_ack  waddr/wdata_commit_id  we_gpr/we_fpr  priv_lvl  debug_mode",
     "commit_stage_i.exception_o   csr_*_commit_*   ex_stage_i.lsu_i.{i_store_unit, i_load_unit}.*"],
    fill=RGBColor(0xE9, 0xEC, 0xF1), border=NAVY, tsize=15, bsize=9.5, mono_body=True)
arrow(s, 6.66, 3.20, 6.66, 3.75, color=BLUE, width=3.0)
label(s, 6.85, 3.22, 5.5, "bind ariane instr_tracer_synth_tap  (resolved in ariane scope; core untouched)",
      size=11, color=BLUE, bold=True, align=PP_ALIGN.LEFT)
tap = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.6), Inches(3.85), Inches(10.1), Inches(2.05))
tap.fill.solid(); tap.fill.fore_color.rgb = RGBColor(0xFD, 0xF3, 0xE3)
tap.line.color.rgb = SIM_B; tap.line.width = Pt(1.5); tap.shadow.inherit = False
label(s, 1.8, 3.95, 9.0, "instr_tracer_synth_tap   ( `ifndef SYNTHESIS  —  simulation wrapper )",
      size=12.5, color=SIM_B, bold=True, align=PP_ALIGN.LEFT)
box(s, 2.05, 4.45, 3.7, 1.20, "instr_tracer_synth",
    ["DUT — synthesizable", "capture + pack + FIFO"], fill=SYNTH, border=SYNTH_B, tsize=13.5)
box(s, 7.55, 4.45, 3.7, 1.20, "instr_tracer_synth_sink",
    ["models the trace consumer", "spike_commit_str() + $fwrite"], fill=SIM, border=SIM_B, tsize=13.5)
arrow(s, 5.75, 4.85, 7.55, 4.85, color=BLUE, width=2.5)
label(s, 5.6, 4.45, 2.1, "trace_valid / beat", size=9.5, color=BLUE, align=PP_ALIGN.CENTER)
arrow(s, 7.55, 5.35, 5.75, 5.35, color=SIM_B, width=1.75)
label(s, 5.6, 5.32, 2.1, "trace_ready", size=9.5, color=SIM_B, align=PP_ALIGN.CENTER)
arrow(s, 9.4, 5.90, 9.4, 6.35, color=GREYTX, width=2.5)
box(s, 7.55, 6.35, 3.7, 0.75, "trace_hart_0_commit.synth.log",
    ["Spike commit-log format"], fill=FILEC, border=FILE_B, tsize=12, bsize=9, shape=MSO_SHAPE.FOLDED_CORNER)
label(s, 1.8, 6.45, 5.6, "(+ trace_hart_0.pkt.hex when EmitPktHex=1, for the real-silicon binary path)",
      size=10.5, color=GREYTX, align=PP_ALIGN.LEFT, italic=True)

# =========================================================================== SLIDE: SIGNAL TAP MAP
s = slide()
title_bar(s, "Where each signal is collected in the CVA6 pipeline",
          "bind resolves these expressions in the `ariane` scope (instr_tracer_synth_bind.sv); "
          "the vector taps reach into Ara at the ara_system level (tb_ara_system_trace.sv).")
_PUR  = RGBColor(0xEE, 0xE6, 0xF7); _PURB = RGBColor(0x78, 0x46, 0xA0)
_ORF  = RGBColor(0xFC, 0xEE, 0xDC); _ORB  = RGBColor(0xC4, 0x76, 0x1C)
_STG  = RGBColor(0xEC, 0xEE, 0xF2)

# ---- top: the CVA6 pipeline + Ara, with numbered tap points ----------------
label(s, 0.40, 1.16, 9.0, "module  ariane   (CVA6 core)  —  bind scope",
      size=12, color=NAVY, bold=True, align=PP_ALIGN.LEFT)
_py, _ph2 = 1.55, 0.98
box(s, 0.40, _py, 1.15, _ph2, "IF · PC", None, fill=_STG, border=NAVY, tsize=11.5)
box(s, 1.70, _py, 1.15, _ph2, "ID",      None, fill=_STG, border=NAVY, tsize=11.5)
box(s, 3.00, _py, 1.50, _ph2, "Issue",   None, fill=_STG, border=NAVY, tsize=11.5)
box(s, 4.65, _py, 2.00, _ph2, "EX",
    ["ALU · MUL · FPU", "lsu_i  ·  csr_regfile"], fill=SYNTH, border=SYNTH_B, tsize=12.5, bsize=8.5, mono_body=True)
box(s, 6.80, _py, 2.45, _ph2, "Commit stage",
    ["commit_ack · exception_o"], fill=SYNTH, border=SYNTH_B, tsize=12.5, bsize=8.5, mono_body=True)
box(s, 9.75, _py, 3.20, _ph2, "Ara vector unit  (i_ara)",
    ["i_dispatcher · gen_lanes[L].i_lane"], fill=GOLD, border=GOLD_B, tsize=12.5, bsize=8.5, mono_body=True)
# pipeline flow + accelerator dispatch
for _x1, _x2 in [(1.55, 1.70), (2.85, 3.00), (4.50, 4.65), (6.65, 6.80)]:
    arrow(s, _x1, _py + _ph2 / 2, _x2, _py + _ph2 / 2, color=GREYTX, width=1.8)
arrow(s, 9.25, _py + _ph2 / 2, 9.75, _py + _ph2 / 2, color=GOLD_B, width=2.2)
label(s, 9.00, 1.30, 1.0, "fu=ACCEL", size=8, color=GOLD_B, bold=True, align=PP_ALIGN.LEFT)
# numbered tap badges on the stage each group observes
badge(s, 4.70, 1.62, 2, color=_PURB,  d=0.34)   # CSR write   -> EX (csr_regfile)
badge(s, 6.22, 1.62, 3, color=_ORB,   d=0.34)   # LSU mem     -> EX (lsu_i)
badge(s, 6.86, 1.62, 1, color=BLUE,   d=0.34)   # identity/result/exception -> Commit
badge(s, 9.80, 1.62, 4, color=GOLD_B, d=0.34)   # vector data -> Ara

# ---- bottom: one callout per tap point, exact hierarchical expressions ------
_cy, _ch = 3.55, 3.18
box(s, 0.35, _cy, 3.05, _ch, "①  COMMIT STAGE",
    ["commit_instr_id_commit → pc·instr",
     "  [p].ex.tval[31:0] → instr word",
     "commit_ack → retired",
     "waddr/wdata_commit_id → rd·wdata",
     "we_gpr/we_fpr_commit_id → we·rd_fpr",
     "priv_lvl · debug_mode → priv·debug",
     "commit_stage_i.exception_o → trap"],
    fill=SYNTH, border=SYNTH_B, tsize=12, bsize=8.5, mono_body=True, align=PP_ALIGN.LEFT)
box(s, 3.50, _cy, 3.05, _ch, "②  CSR WRITE   “c<a>_<n> 0x..”",
    ["csr_commit_commit_ex → pulse",
     "csr_op_commit_csr → W/S/C",
     "csr_addr_ex_csr → csr_addr",
     "csr_wdata_commit_csr → operand",
     "csr_rdata_csr_commit → old val",
     "",
     "pre-WARL → --ignore-csr-val"],
    fill=_PUR, border=_PURB, tsize=11, bsize=8.5, mono_body=True, align=PP_ALIGN.LEFT)
box(s, 6.65, _cy, 3.05, _ch, "③  LSU MEM   “mem 0x.. 0x..”",
    ["ex_stage_i.lsu_i :",
     "store_buffer_i.valid/paddr/",
     "   data/data_size_i",
     "load_unit.req_port_o.tag_valid/",
     "   kill_req/data_size",
     "load_unit.paddr_i",
     "flush_ctrl_ex → flush_addr_i"],
    fill=_ORF, border=_ORB, tsize=11, bsize=8.5, mono_body=True, align=PP_ALIGN.LEFT)
box(s, 9.80, _cy, 3.18, _ch, "④  ARA VECTOR   “v<vd> 0x..”",
    ["i_dispatcher.vtype_q.vsew/",
     "   vlmul · vl_q → vsew·vl",
     "gen_lanes[L].i_lane.vrf_wen/",
     "   addr/wdata/be",
     "   → shadow VRF → v[vd] data",
     "",
     "system-level taps (not bind)"],
    fill=GOLD, border=GOLD_B, tsize=11, bsize=8.5, mono_body=True, align=PP_ALIGN.LEFT)
label(s, 0.35, 6.86, 12.6,
      "Scalar taps: instr_tracer_synth_bind.sv (bind ariane).   "
      "Vector taps reach into i_ara, wired at ara_system level (tb_ara_system_trace.sv / exampleinit.sv).",
      size=9.5, color=GREYTX, align=PP_ALIGN.LEFT, italic=True)

# =========================================================================== SLIDE 5
s = slide()
title_bar(s, "Co-simulation verification flow",
          "Prove the synthesizable tracer reproduces the same commit log as the golden reference — instruction for instruction.")
box(s, 4.55, 1.30, 4.25, 1.10, "RTL simulation",
    ["make app=… simc  (QuestaSim)  /  verilate + simv (Verilator)",
     "ara_system + CVA6 + bound tracer"], fill=RGBColor(0xE9, 0xEC, 0xF1), border=NAVY, tsize=14, bsize=9.5)
box(s, 0.55, 1.45, 3.4, 0.8, "app.elf / riscv-dv  program.hex",
    fill=FILEC, border=FILE_B, tsize=12, shape=MSO_SHAPE.FOLDED_CORNER)
arrow(s, 3.95, 1.85, 4.55, 1.85, color=BLUE, width=2.25)
box(s, 3.05, 3.05, 3.5, 1.05, "trace_hart_0_commit.log",
    ["GOLDEN — original sim tracer", "(QuestaSim only)"],
    fill=GOLD, border=GOLD_B, tsize=12.5, bsize=9.5, shape=MSO_SHAPE.FOLDED_CORNER)
box(s, 6.85, 3.05, 3.5, 1.05, "trace_hart_0_commit.synth.log",
    ["DUT output — synth tracer sink", "(Spike commit-log format)"],
    fill=SYNTH, border=SYNTH_B, tsize=12.5, bsize=9.5, shape=MSO_SHAPE.FOLDED_CORNER)
arrow(s, 5.55, 2.40, 4.40, 3.05, color=GOLD_B)
arrow(s, 7.20, 2.40, 8.40, 3.05, color=SYNTH_B)
box(s, 5.15, 4.55, 3.0, 0.85, "diff", ["byte-for-byte equal"],
    fill=RGBColor(0xEF, 0xEF, 0xEF), border=GREYTX, tsize=15, bsize=9.5)
arrow(s, 4.70, 4.10, 6.10, 4.55, color=GOLD_B)
arrow(s, 8.60, 4.10, 7.20, 4.55, color=SYNTH_B)
box(s, 5.55, 5.65, 2.2, 0.65, "PASS / FAIL", fill=GOLD, border=GOLD_B, tsize=14)
arrow(s, 6.65, 5.40, 6.65, 5.65, color=GREYTX)
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.0), Inches(2.55), Inches(0.02), Inches(4.0))
line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
line.line.fill.background(); line.shadow.inherit = False
label(s, 9.2, 2.55, 4.0, "Alternative golden references", size=13, color=NAVY, bold=True, align=PP_ALIGN.LEFT)
box(s, 9.25, 3.05, 3.85, 0.9, "Spike ISS co-simulation",
    ["spike --log-commits app.elf → spike.log"], fill=GOLD, border=GOLD_B, tsize=12.5, bsize=9, mono_body=True)
arrow(s, 11.15, 3.95, 11.15, 4.35, color=GOLD_B)
box(s, 9.25, 4.35, 3.85, 1.0, "compare_spike_log.py",
    ["drops disasm / boot-ROM, normalizes", "regs & CSRs (--ignore-csr-val …)"],
    fill=SIM, border=SIM_B, tsize=12.5, bsize=9)
arrow(s, 11.15, 5.35, 11.15, 5.70, color=GREYTX)
box(s, 9.9, 5.70, 2.5, 0.6, "PASS / FAIL", fill=GOLD, border=GOLD_B, tsize=13)
label(s, 9.25, 6.55, 3.9, "Silicon path: pkt.hex → spike_trace_decode.py → .log → diff",
      size=9.5, color=GREYTX, align=PP_ALIGN.LEFT, italic=True)

# =========================================================================== SLIDE 6
s = slide()
title_bar(s, "Files, run commands & caveats")
label(s, 0.6, 1.15, 6.3, "Source files", size=16, color=NAVY, bold=True, align=PP_ALIGN.LEFT)
files = [
    ("instr_tracer_synth_pkg.sv", "binary record format (commit_log_pkt_t / beat)"),
    ("instr_tracer_synth.sv", "the tracer: capture, pack, FIFO, ready/valid out"),
    ("instr_tracer_addr_fifo.sv", "bounded store/load address-tracking FIFO"),
    ("instr_tracer_synth_sink.sv", "sim sink: prints Spike log directly in SV"),
    ("instr_tracer_synth_tap.sv", "sim wrapper: tracer + sink"),
    ("instr_tracer_synth_bind.sv", "non-invasive bind into `ariane`"),
    ("tb_ara_system_trace.sv", "standalone ara_system + tracer testbench"),
    ("scripts/spike_trace_decode.py", "off-chip binary → Spike ASCII decoder"),
    ("scripts/compare_spike_log.py", "compare against spike --log-commits"),
]
y = 1.6
for f, d in files:
    label(s, 0.7, y, 4.0, f, size=11, color=BLUE, bold=True, align=PP_ALIGN.LEFT, h=0.28)
    label(s, 0.9, y + 0.215, 6.0, d, size=9.5, color=GREYTX, align=PP_ALIGN.LEFT, h=0.25)
    y += 0.52
label(s, 7.1, 1.15, 5.8, "Run (fast path, machine already set up)", size=16, color=NAVY, bold=True, align=PP_ALIGN.LEFT)
cmdbox = box(s, 7.1, 1.6, 5.65, 1.9, "", None, fill=RGBColor(0x1B, 0x20, 0x28),
             border=RGBColor(0x1B, 0x20, 0x28), shape=MSO_SHAPE.RECTANGLE)
tf = cmdbox.text_frame; tf.vertical_anchor = MSO_ANCHOR.TOP
tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1)
for i, c in enumerate(["make -C apps bin/hello_world",
                       "make -C hardware app=hello_world simc",
                       "cd hardware/build",
                       "diff trace_hart_0_commit.log \\",
                       "     trace_hart_0_commit.synth.log"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = c
    r.font.size = Pt(11); r.font.name = MONO; r.font.color.rgb = RGBColor(0x8C, 0xE0, 0x9A)
label(s, 7.1, 3.75, 5.8, "Caveats (non-invasive tracer)", size=16, color=NAVY, bold=True, align=PP_ALIGN.LEFT)
yy = 4.2
for n in [
    "CSR write value is reconstructed PRE-WARL mask → WARL CSRs may differ "
    "from Spike's post-WARL value (use --ignore-csr-val).",
    "Store data is un-rotated (byte_ror64) to undo CVA6 data_align before slicing.",
    "Exceptions/interrupts always reported on port 0; debug-mode BREAKPOINT suppressed.",
    "overflow_o pulses if the trace FIFO (or an addr FIFO) over/underflows → "
    "trace no longer gap-free; raise FifoDepth.",
]:
    bl = s.shapes.add_textbox(Inches(7.1), Inches(yy), Inches(5.7), Inches(0.7))
    tf = bl.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r0 = p.add_run(); r0.text = "▸ "; r0.font.size = Pt(10.5); r0.font.color.rgb = BLUE; r0.font.bold = True
    r = p.add_run(); r.text = n; r.font.size = Pt(10.5); r.font.color.rgb = GREYTX
    yy += 0.72

# =========================================================================== save
out_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "docs")
os.makedirs(out_dir, exist_ok=True)
out = os.path.join(out_dir, "instr_tracer_synth_arch.pptx")
prs.save(out)
print("wrote", out, "(%d slides)" % len(prs.slides._sldIdLst))
